"""
avatarpipeline.narration.slide_renderer — Render PPTX slides to PNG images.

Strategy (in priority order):
  1. LibreOffice headless — highest fidelity, handles charts/images/themes.
  2. python-pptx + PIL fallback — works offline without LibreOffice; renders
     background, title and body text in a clean layout.
"""
from __future__ import annotations

import shutil
import subprocess
from pathlib import Path

from loguru import logger

# ── LibreOffice discovery ────────────────────────────────────────────────────

_SOFFICE_CANDIDATES: list[str] = [
    "soffice",
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    "/opt/homebrew/bin/soffice",
    "/usr/local/bin/soffice",
    "/opt/local/bin/soffice",
    "/usr/bin/soffice",
]


def _find_soffice() -> str | None:
    for candidate in _SOFFICE_CANDIDATES:
        found = shutil.which(candidate)
        if found:
            return found
        if Path(candidate).exists():
            return candidate
    return None


# ── LibreOffice renderer ─────────────────────────────────────────────────────

def _render_libreoffice(pptx_path: Path, output_dir: Path) -> list[Path]:
    soffice = _find_soffice()
    if not soffice:
        raise FileNotFoundError("LibreOffice not found in any standard path.")

    output_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        soffice, "--headless", "--norestore",
        "--convert-to", "png",
        "--outdir", str(output_dir),
        str(pptx_path),
    ]
    r = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if r.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed:\n{r.stderr[-600:]}")

    stem = pptx_path.stem
    slides = sorted(output_dir.glob(f"{stem}*.png"))
    if not slides:
        raise RuntimeError("LibreOffice produced no PNG output.")
    return slides


# ── PIL-based fallback renderer ──────────────────────────────────────────────

_FONT_CANDIDATES: list[str] = [
    "/System/Library/Fonts/Helvetica.ttc",
    "/System/Library/Fonts/Arial.ttf",
    "/System/Library/Fonts/SFNSDisplay.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
]


def _load_font(size: int):
    from PIL import ImageFont
    for path in _FONT_CANDIDATES:
        if Path(path).exists():
            try:
                return ImageFont.truetype(path, size=size)
            except Exception:
                continue
    return ImageFont.load_default()


def _get_bg_color(slide) -> tuple[int, int, int]:
    """Extract slide background as an (R, G, B) tuple; returns dark-navy fallback."""
    try:
        fill = slide.background.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc.type is not None:
                rgb = fc.rgb
                return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    return (22, 40, 72)


def _draw_wrapped(
    draw, text: str, x: int, y: int, max_w: int, font, fill
) -> int:
    """Wrap ``text`` at ``max_w`` pixels and draw it; returns updated y."""
    words = text.split()
    if not words:
        return y
    lines: list[str] = []
    current: list[str] = []
    for word in words:
        candidate = " ".join(current + [word])
        bbox = draw.textbbox((0, 0), candidate, font=font)
        if bbox[2] - bbox[0] <= max_w:
            current.append(word)
        else:
            if current:
                lines.append(" ".join(current))
            current = [word]
    if current:
        lines.append(" ".join(current))

    for line in lines:
        draw.text((x, y), line, fill=fill, font=font)
        bbox = draw.textbbox((0, 0), line, font=font)
        y += int((bbox[3] - bbox[1]) * 1.32)
    return y


def _render_slide_pil(slide, prs, width: int, height: int, slide_num: int, total: int):
    from PIL import Image, ImageDraw

    bg = _get_bg_color(slide)
    img = Image.new("RGB", (width, height), bg)
    draw = ImageDraw.Draw(img)

    # ── Header bar ──────────────────────────────────────────────────────────
    bar_h = max(36, height // 28)
    bar_dark = tuple(max(c - 18, 0) for c in bg)
    draw.rectangle([(0, 0), (width, bar_h)], fill=bar_dark)

    font_sm = _load_font(max(14, int(height * 0.018)))
    counter = f"{slide_num} / {total}"
    bbox_c = draw.textbbox((0, 0), counter, font=font_sm)
    draw.text(
        (width - (bbox_c[2] - bbox_c[0]) - 20, (bar_h - (bbox_c[3] - bbox_c[1])) // 2),
        counter,
        fill=(160, 180, 210),
        font=font_sm,
    )

    # ── Collect title and body text ──────────────────────────────────────────
    title_text: str | None = None
    body_texts: list[str] = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        ph = getattr(shape, "placeholder_format", None)
        ph_idx = ph.idx if ph is not None else -1
        if ph_idx == 0:
            title_text = text
        else:
            body_texts.append(text)

    pad_x = int(width * 0.06)
    max_w = int(width * 0.88)
    y = bar_h + int(height * 0.07)

    # ── Title ────────────────────────────────────────────────────────────────
    if title_text:
        font_title = _load_font(max(24, int(height * 0.055)))
        y = _draw_wrapped(draw, title_text, pad_x, y, max_w, font_title, (255, 255, 255))
        y += int(height * 0.03)
        # Divider
        div_col = tuple(min(c + 55, 255) for c in bg)
        draw.line([(pad_x, y), (pad_x + max_w, y)], fill=div_col, width=2)
        y += int(height * 0.04)

    # ── Body ─────────────────────────────────────────────────────────────────
    font_body = _load_font(max(16, int(height * 0.032)))
    for text in body_texts:
        if y > height * 0.88:
            break
        y = _draw_wrapped(draw, f"• {text}", pad_x, y, max_w, font_body, (205, 215, 235))
        y += int(height * 0.022)

    return img


def _render_pil(pptx_path: Path, output_dir: Path, render_width: int = 1920) -> list[Path]:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    output_dir.mkdir(parents=True, exist_ok=True)

    emu_w = prs.slide_width or 9_144_000
    emu_h = prs.slide_height or 5_143_500
    render_height = int(render_width * emu_h / emu_w)
    # Ensure even pixel dimensions for ffmpeg
    render_height += render_height % 2
    render_width += render_width % 2

    paths: list[Path] = []
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        img = _render_slide_pil(slide, prs, render_width, render_height, i, total)
        out = output_dir / f"slide_{i:03d}.png"
        img.save(str(out))
        logger.debug(f"Rendered slide {i}/{total} via PIL")
        paths.append(out)
    return paths


# ── Public entry point ───────────────────────────────────────────────────────

def render_slides(pptx_path: str | Path, output_dir: str | Path) -> list[Path]:
    """Render all PPTX slides to PNG images.

    Tries LibreOffice first (best fidelity).  Falls back to a python-pptx +
    PIL renderer that handles background colours and text shapes.

    Args:
        pptx_path:  Path to the ``.pptx`` file.
        output_dir: Directory where ``slide_NNN.png`` files are written.

    Returns:
        List of :class:`~pathlib.Path` objects in slide order.
    """
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)

    if _find_soffice():
        logger.info("Rendering slides via LibreOffice")
        try:
            return _render_libreoffice(pptx_path, output_dir)
        except Exception as exc:
            logger.warning(f"LibreOffice render failed ({exc}); falling back to PIL")

    logger.info("Rendering slides via python-pptx + PIL")
    return _render_pil(pptx_path, output_dir)
