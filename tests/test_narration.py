#!/usr/bin/env python3
"""
tests.test_narration — Focused tests for the slide narrator pipeline.
"""

import subprocess
import sys
from pathlib import Path

import numpy as np
import soundfile as sf
from PIL import Image
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))


def _make_pptx(path: Path, slide_count: int) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[6]
    for _ in range(slide_count):
        prs.slides.add_slide(layout)
    prs.save(str(path))


def _write_silence(path: Path, seconds: float, sample_rate: int = 16000) -> None:
    samples = max(1, int(seconds * sample_rate))
    audio = np.zeros(samples, dtype=np.float32)
    sf.write(path, audio, sample_rate)


def _probe_duration(path: Path) -> float:
    cmd = [
        "ffprobe", "-v", "error",
        "-show_entries", "format=duration",
        "-of", "default=noprint_wrappers=1:nokey=1",
        str(path),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, check=True)
    return float(result.stdout.strip())


def test_validate_sync_accepts_flexible_json_schema(tmp_path):
    from avatarpipeline.narration.validator import validate_sync

    pptx_path = tmp_path / "deck.pptx"
    _make_pptx(pptx_path, 2)

    json_data = [
        {"text": "Intro narration", "duration": 4},
        {"script": "Second narration", "pause": 0.5},
    ]

    result = validate_sync(pptx_path, json_data)

    assert result.ok
    assert result.slide_count == 2
    assert result.json_count == 2
    assert result.json_data["slides"][0]["slide_number"] == 1
    assert result.json_data["slides"][0]["display_seconds"] == 4.0
    assert result.json_data["slides"][1]["pause_seconds"] == 0.5
    assert any("JSON order was used" in warning for warning in result.warnings)


def test_validate_sync_accepts_slide_number_keyed_dict(tmp_path):
    from avatarpipeline.narration.validator import validate_sync

    pptx_path = tmp_path / "deck.pptx"
    _make_pptx(pptx_path, 2)

    json_data = {
        "slides": {
            "1": {"narration": "First"},
            "2": {"voiceover": "Second", "duration_seconds": 2.5},
        }
    }

    result = validate_sync(pptx_path, json_data)

    assert result.ok
    assert [entry["slide_number"] for entry in result.json_data["slides"]] == [1, 2]
    assert result.json_data["slides"][1]["display_seconds"] == 2.5


def test_compose_narrated_video_generates_audio_before_render_and_uses_slide_timing(tmp_path, monkeypatch):
    import avatarpipeline.narration.composer as composer
    import avatarpipeline.voice.kokoro as kokoro_mod

    pptx_path = tmp_path / "deck.pptx"
    _make_pptx(pptx_path, 2)

    slide_images = []
    for idx, color in enumerate(((30, 64, 175), (15, 118, 110)), start=1):
        img_path = tmp_path / f"slide_{idx:03d}.png"
        Image.new("RGB", (1920, 1080), color).save(img_path)
        slide_images.append(img_path)

    def fake_render_slides(_pptx_path, _output_dir):
        return slide_images

    class FakeVoiceGenerator:
        def generate(self, text, voice=None, out_path=None):
            duration = 1.0 if "First" in text else 1.0
            _write_silence(Path(out_path), duration)
            return str(out_path)

    monkeypatch.setattr(composer, "render_slides", fake_render_slides)
    monkeypatch.setattr(kokoro_mod, "VoiceGenerator", FakeVoiceGenerator)

    json_data = {
        "default_pause_seconds": 0.25,
        "slides": [
            {"slide_number": 1, "narration": "First slide narration", "duration_seconds": 3.0},
            {"slide_number": 2, "narration": "Second slide narration", "pause_seconds": 0.5},
        ],
    }

    events = list(
        composer.compose_narrated_video(
            pptx_path=pptx_path,
            json_data=json_data,
            output_path=tmp_path / "narrated.mp4",
            voice="af_heart",
            pause_seconds=1.0,
        )
    )

    messages = [message for message, _ in events]
    render_index = next(i for i, message in enumerate(messages) if message.startswith("Rendering slides"))
    tts_index = next(i for i, message in enumerate(messages) if message.startswith("TTS"))
    output_path = Path(events[-1][1])

    assert tts_index < render_index
    assert output_path.exists()

    duration = _probe_duration(output_path)
    assert 4.4 <= duration <= 5.2
